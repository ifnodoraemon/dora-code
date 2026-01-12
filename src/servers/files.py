import os
import glob
from pathlib import Path
from typing import List, Optional
from mcp.server.fastmcp import FastMCP

# Document Parsers
import pdfplumber
import docx
from pptx import Presentation
import openpyxl

# Vision
import google.generativeai as genai
from PIL import Image

mcp = FastMCP("PolymathFilesystem")

# --------------------------
# Vision Helper
# --------------------------
def ocr_with_gemini(image_path: str) -> str:
    """Fallback OCR using Gemini if text extraction fails."""
    api_key = os.getenv("GOOGLE_API_KEY")
    if not api_key:
        return "[Error: GOOGLE_API_KEY needed for OCR]"
    
    try:
        genai.configure(api_key=api_key)
        model = genai.GenerativeModel("gemini-1.5-flash")
        img = Image.open(image_path)
        response = model.generate_content(["Extract all text from this page.", img])
        return response.text
    except Exception as e:
        return f"[OCR Failed: {e}]"

# --------------------------
# Parsers
# --------------------------
def parse_pdf(path: str) -> str:
    text = []
    with pdfplumber.open(path) as pdf:
        for i, page in enumerate(pdf.pages):
            page_text = page.extract_text()
            # 如果提取为空，尝试 OCR (只尝试第一页以节省 token，或者按需)
            if not page_text or len(page_text.strip()) < 10:
                # 将 PDF 页面转为图片进行 OCR (需要 pdf2image，这里简化处理)
                # 由于环境限制，我们先假设用户会手动把扫描件转图
                # 或者提示用户："此 PDF 似乎是扫描件，请先转为图片"
                page_text = "[Warning: Scanned PDF detected. Text extraction unreliable.]"
            text.append(page_text)
    return "\n".join(text)

def parse_image(path: str) -> str:
    """直接处理图片文件"""
    return ocr_with_gemini(path)

def parse_docx(path: str) -> str:
    doc = docx.Document(path)
    return "\n".join([p.text for p in doc.paragraphs])

def parse_pptx(path: str) -> str:
    prs = Presentation(path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def parse_xlsx(path: str) -> str:
    wb = openpyxl.load_workbook(path, data_only=True)
    text = []
    for sheet in wb.sheetnames:
        text.append(f"--- Sheet: {sheet} ---")
        ws = wb[sheet]
        for row in ws.iter_rows(values_only=True):
            text.append("\t".join([str(c) for c in row if c is not None]))
    return "\n".join(text)

# --------------------------
# Tools
# --------------------------
def validate_path(path: str) -> str:
    """Validate that path is within the current working directory."""
    abs_path = os.path.abspath(path)
    base_dir = os.getcwd()
    if not abs_path.startswith(base_dir):
        raise PermissionError(f"Access Denied: Path {path} is outside of the workspace sandbox.")
    return abs_path

@mcp.tool()
def list_directory(path: str = ".") -> str:
    """List files and directories at the given path."""
    valid_path = validate_path(path)
    if not os.path.exists(valid_path): return "Error: Path not found."
    try:
        return "\n".join(os.listdir(valid_path))
    except Exception as e:
        return f"Error: {e}"

@mcp.tool()
def read_file(path: str) -> str:
    """
    Intelligently read a file. Supports PDF, Office, and Images (via OCR).
    """
    valid_path = validate_path(path)
    if not os.path.exists(valid_path): return "Error: File not found."
    
    ext = os.path.splitext(path)[1].lower()
    
    try:
        if ext == ".pdf":
            return parse_pdf(valid_path)
        elif ext == ".docx":
            return parse_docx(valid_path)
        elif ext == ".pptx":
            return parse_pptx(valid_path)
        elif ext in [".xlsx", ".xls"]:
            return parse_xlsx(valid_path)
        elif ext in [".png", ".jpg", ".jpeg", ".webp"]:
            return parse_image(valid_path)
        else:
            with open(valid_path, "r", encoding="utf-8", errors="replace") as f:
                return f.read()
    except Exception as e:
        return f"Error reading file: {str(e)}"

@mcp.tool()
def write_file(path: str, content: str) -> str:
    """Write text content to a file."""
    valid_path = validate_path(path)
    try:
        os.makedirs(os.path.dirname(valid_path), exist_ok=True)
        with open(valid_path, "w", encoding="utf-8") as f:
            f.write(content)
        return f"Successfully wrote to {path}"
    except Exception as e:
        return f"Error writing file: {str(e)}"

if __name__ == "__main__":
    mcp.run()