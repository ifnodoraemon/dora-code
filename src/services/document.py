import pdfplumber
import docx
from pptx import Presentation
import openpyxl

def parse_pdf(path: str) -> str:
    text = []
    try:
        with pdfplumber.open(path) as pdf:
            for i, page in enumerate(pdf.pages):
                page_text = page.extract_text()
                if not page_text or len(page_text.strip()) < 20:
                    page_text = f"[Page {i+1} is likely scanned. OCR required.]"
                text.append(page_text)
        return "\n".join(text)
    except Exception as e:
        return f"[PDF Error: {e}]"

def parse_docx(path: str) -> str:
    try:
        doc = docx.Document(path)
        return "\n".join([p.text for p in doc.paragraphs])
    except Exception as e:
        return f"[DOCX Error: {e}]"

def parse_pptx(path: str) -> str:
    try:
        prs = Presentation(path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    except Exception as e:
        return f"[PPTX Error: {e}]"

def parse_xlsx(path: str) -> str:
    try:
        wb = openpyxl.load_workbook(path, data_only=True)
        text = []
        for sheet in wb.sheetnames:
            text.append(f"--- Sheet: {sheet} ---")
            ws = wb[sheet]
            for row in ws.iter_rows(values_only=True):
                text.append("\t".join([str(c) for c in row if c is not None]))
        return "\n".join(text)
    except Exception as e:
        return f"[XLSX Error: {e}]"
